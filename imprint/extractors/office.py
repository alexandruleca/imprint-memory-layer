"""PowerPoint (.pptx), Excel (.xlsx), CSV / TSV extractors.

PPTX: python-pptx — concatenate slide text frames.
XLSX: openpyxl — flatten each sheet to "header: value" rows to preserve
      column semantics for retrieval.
CSV/TSV: stdlib csv — first row treated as header if non-numeric.
"""

from __future__ import annotations

import csv as _csv
import io
import os

from . import ExtractedDoc, ExtractorUnavailable, ExtractionError, register_ext, register_mime


# ── PPTX ───────────────────────────────────────────────────────
def _extract_pptx_bytes(data: bytes, source_url: str = "") -> ExtractedDoc:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as e:
        raise ExtractorUnavailable("python-pptx not installed — `pip install python-pptx`") from e

    try:
        pres = Presentation(io.BytesIO(data))
    except Exception as e:
        raise ExtractionError(f"pptx parse failed: {e}") from e

    slides = []
    for i, slide in enumerate(pres.slides, 1):
        bits = [f"## Slide {i}"]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    bits.append(line)
        if slide.has_notes_slide:
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    bits.append(f"Notes: {notes}")
            except Exception:
                pass
        if len(bits) > 1:
            slides.append("\n".join(bits))

    text = "\n\n".join(slides).strip()
    return ExtractedDoc(
        text=text,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        metadata={"slide_count": len(pres.slides)},
        chunk_mode="prose",
    )


def _extract_pptx(path: str) -> ExtractedDoc:
    with open(path, "rb") as f:
        doc = _extract_pptx_bytes(f.read())
    doc.metadata.setdefault("filename", os.path.basename(path))
    return doc


register_ext(".pptx", _extract_pptx)
register_mime(
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    _extract_pptx_bytes,
)


# ── XLSX ───────────────────────────────────────────────────────
def _extract_xlsx_bytes(data: bytes, source_url: str = "") -> ExtractedDoc:
    try:
        import openpyxl  # type: ignore
    except ImportError as e:
        raise ExtractorUnavailable("openpyxl not installed — `pip install openpyxl`") from e

    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as e:
        raise ExtractionError(f"xlsx parse failed: {e}") from e

    sheets = []
    for sheet in wb.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
        bits = [f"## Sheet: {sheet.title}"]
        headers = [str(h) if h is not None else "" for h in rows[0]]
        has_header = any(h and not str(h).replace(".", "").replace("-", "").isdigit() for h in headers)
        if has_header:
            for row in rows[1:]:
                pairs = []
                for h, v in zip(headers, row):
                    if v is None or v == "":
                        continue
                    pairs.append(f"{h}: {v}" if h else str(v))
                if pairs:
                    bits.append(" | ".join(pairs))
        else:
            for row in rows:
                cells = [str(c) for c in row if c is not None and c != ""]
                if cells:
                    bits.append(" | ".join(cells))
        if len(bits) > 1:
            sheets.append("\n".join(bits))

    text = "\n\n".join(sheets).strip()
    return ExtractedDoc(
        text=text,
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        metadata={"sheet_count": len(wb.worksheets)},
        chunk_mode="prose",
    )


def _extract_xlsx(path: str) -> ExtractedDoc:
    with open(path, "rb") as f:
        doc = _extract_xlsx_bytes(f.read())
    doc.metadata.setdefault("filename", os.path.basename(path))
    return doc


register_ext(".xlsx", _extract_xlsx)
register_mime(
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    _extract_xlsx_bytes,
)


# ── CSV / TSV (stdlib, always available) ───────────────────────
def _extract_csv_text(text: str, delimiter: str) -> str:
    reader = _csv.reader(io.StringIO(text), delimiter=delimiter)
    rows = list(reader)
    if not rows:
        return ""
    headers = rows[0]
    has_header = any(h and not h.replace(".", "").replace("-", "").isdigit() for h in headers)
    out: list[str] = []
    if has_header:
        for row in rows[1:]:
            pairs = [f"{h}: {v}" for h, v in zip(headers, row) if v]
            if pairs:
                out.append(" | ".join(pairs))
    else:
        for row in rows:
            cells = [c for c in row if c]
            if cells:
                out.append(" | ".join(cells))
    return "\n".join(out)


def _extract_csv(path: str) -> ExtractedDoc:
    with open(path, "r", errors="ignore", newline="") as f:
        raw = f.read()
    ext = os.path.splitext(path)[1].lower()
    delim = "\t" if ext == ".tsv" else ","
    try:
        text = _extract_csv_text(raw, delim)
    except Exception as e:
        raise ExtractionError(f"csv parse failed: {e}") from e
    return ExtractedDoc(
        text=text,
        mime="text/tab-separated-values" if ext == ".tsv" else "text/csv",
        metadata={"filename": os.path.basename(path)},
        chunk_mode="prose",
    )


def _extract_csv_bytes(data: bytes, source_url: str = "") -> ExtractedDoc:
    raw = data.decode("utf-8", errors="ignore")
    text = _extract_csv_text(raw, ",")
    return ExtractedDoc(text=text, mime="text/csv", metadata={}, chunk_mode="prose")


register_ext(".csv", _extract_csv)
register_ext(".tsv", _extract_csv)
register_mime("text/csv", _extract_csv_bytes)
register_mime("text/tab-separated-values", _extract_csv_bytes)
